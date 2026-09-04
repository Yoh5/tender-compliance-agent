"""Build `docs/deck.pptx`, the slides for the five-minute submission video.

    python scripts/make_deck.py

Kept in the repository rather than run and deleted, for the same reason as the
other scripts here: a deck edited by hand drifts from the script it illustrates,
and nobody notices until the take. This one is generated from the same facts as
`docs/video.md`, so a figure that changes there is changed in one place.

THE SLIDES ARE NOT THE VIDEO. Beats 6 and 7 of `docs/video.md` are a screen
recording of two live runs; the deck carries a marker for them so the running
order stays intact while rehearsing, and those two slides are cut away in the
edit. Everything a slide claims is measured — the numbers here are the ones
observed on 2026-09-03/04 and named as such, because a slide is read as a
statement of fact in a way a spoken sentence is not.

The palette is the report's own (`report.py`), so the deck and the thing it
describes look like one product rather than two.
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SORTIE = ROOT / "docs" / "deck.pptx"
SCHEMA = ROOT / "docs" / "architecture.svg"

# --------------------------------------------------------------- the palette
# Lifted from report.py so the deck and the report are the same object.
PAPER = RGBColor(0xF6, 0xF8, 0xFB)
INK = RGBColor(0x14, 0x20, 0x3A)      # the navy of French official print
MUTED = RGBColor(0x5B, 0x68, 0x80)
RULE = RGBColor(0xCC, 0xD6, 0xE4)
STAMP = RGBColor(0xA8, 0x24, 0x3B)    # madder; blockers only, nothing else
GOOD = RGBColor(0x2C, 0x5F, 0x4E)
REVIEW = RGBColor(0x7D, 0x5C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xC9, 0xD6, 0xE8)

SANS = "Segoe UI"
SERIF = "Georgia"

LARGEUR, HAUTEUR = 13.333, 7.5
MARGE = 0.95
COLONNE = LARGEUR - 2 * MARGE


# The diagram's own heading lives above y=88 in its viewBox. It is cropped out
# for the deck: the slide already carries a title, and printing two is the kind
# of thing that reads as assembled rather than designed.
SCHEMA_HAUT = 88

def _rendre_le_schema() -> tuple[Path, float]:
    """architecture.svg → PNG, because PowerPoint will not place an SVG.

    Rendered rather than kept as a second file: two copies of one diagram drift,
    and the one that drifts is always the one nobody opens. Returns the aspect
    ratio too, so the caller sizes the picture from the image instead of from a
    number typed by hand — the first version overflowed the slide and cut the
    last line of the diagram off.
    """
    import fitz

    cible = ROOT / "docs" / "architecture.png"
    document = fitz.open("svg", SCHEMA.read_bytes())
    page = document[0]
    zone = fitz.Rect(0, SCHEMA_HAUT, page.rect.width, page.rect.height)
    pixels = page.get_pixmap(matrix=fitz.Matrix(2.6, 2.6), clip=zone, alpha=False)
    pixels.save(cible)
    return cible, pixels.width / pixels.height


# ------------------------------------------------------------------ plumbing

def _texte(diapo, gauche, haut, largeur, hauteur):
    boite = diapo.shapes.add_textbox(Inches(gauche), Inches(haut),
                                     Inches(largeur), Inches(hauteur))
    cadre = boite.text_frame
    cadre.word_wrap = True
    cadre.margin_left = cadre.margin_right = 0
    cadre.margin_top = cadre.margin_bottom = 0
    return cadre


def _ligne(cadre, texte, *, taille, couleur, gras=False, police=SANS,
           interligne=1.15, apres=0, premiere=False, aligne=PP_ALIGN.LEFT,
           espacement=None):
    para = cadre.paragraphs[0] if premiere else cadre.add_paragraph()
    para.alignment = aligne
    para.line_spacing = interligne
    para.space_after = Pt(apres)
    passage = para.add_run()
    passage.text = texte
    passage.font.size = Pt(taille)
    passage.font.name = police
    passage.font.bold = gras
    passage.font.color.rgb = couleur
    if espacement is not None:
        # Letter-spacing. `spc` is a plain attribute on a:rPr, in hundredths of
        # a point, and it is not namespaced — qn() refuses it.
        passage.font._rPr.set("spc", str(int(espacement * 100)))
    return para


def _filet(diapo, gauche, haut, largeur=1.15, couleur=STAMP, epaisseur=0.055):
    """The recurring mark: a short madder rule above every title.

    One motif, used the same way on every slide, is what makes a deck read as
    designed rather than assembled.
    """
    forme = diapo.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(gauche),
                                   Inches(haut), Inches(largeur),
                                   Inches(epaisseur))
    forme.fill.solid()
    forme.fill.fore_color.rgb = couleur
    forme.line.fill.background()
    forme.shadow.inherit = False
    return forme


def _fond(diapo, couleur):
    diapo.background.fill.solid()
    diapo.background.fill.fore_color.rgb = couleur


def _notes(diapo, texte):
    diapo.notes_slide.notes_text_frame.text = texte.strip()


def _page(prs, couleur=PAPER):
    diapo = prs.slides.add_slide(prs.slide_layouts[6])
    _fond(diapo, couleur)
    return diapo


def _titre(diapo, titre, *, haut=1.05, sur_fond_sombre=False):
    _filet(diapo, MARGE, haut, couleur=STAMP)
    cadre = _texte(diapo, MARGE, haut + 0.28, COLONNE, 1.0)
    _ligne(cadre, titre, taille=34, gras=True,
           couleur=WHITE if sur_fond_sombre else INK, premiere=True)
    return haut + 1.25


# -------------------------------------------------------------------- slides

def construire() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(LARGEUR), Inches(HAUTEUR)

    # 1 ─ title ------------------------------------------------------------
    d = _page(prs, INK)
    _filet(d, MARGE, 2.5, largeur=1.6)
    c = _texte(d, MARGE, 2.9, COLONNE, 2.4)
    _ligne(c, "Tender Compliance Agent", taille=54, gras=True, couleur=WHITE,
           premiere=True, apres=14)
    _ligne(c, "Public bids are rejected on paperwork before anyone reads them.",
           taille=23, couleur=PALE, police=SERIF, apres=6)
    _ligne(c, "This agent finds the gaps first.", taille=23, couleur=PALE,
           police=SERIF)
    c = _texte(d, MARGE, 6.05, COLONNE, 0.7)
    _ligne(c, "AGENTS FOR HUMANS  ·  PROFESSIONAL AGENTS TRACK  ·  BUILT WITH "
              "THE STRANDS AGENTS SDK",
           taille=11, couleur=MUTED, premiere=True, espacement=1.6)
    _notes(d, """
Twenty seconds. Say what it is before saying why it should exist.

"This is a compliance agent for public tenders. You give it the buyer's
consultation file and the list of papers your company holds, and it tells you
what is missing, what expires too soon, and what a human still has to look at."
""")

    # 2 ─ the problem ------------------------------------------------------
    d = _page(prs)
    bas = _titre(d, "A bid dies on the paperwork")
    c = _texte(d, MARGE, bas, 5.6, 3.0)
    _ligne(c, "€100 bn", taille=46, gras=True, couleur=INK, premiere=True,
           apres=10)
    _ligne(c, "of French public contracts a year, roughly — and a bid can be "
              "thrown out "
              "before anyone reads it. Not on price. Not on technique. Because "
              "an attestation expired three days before the deadline, or a form "
              "was missing.",
           taille=15, couleur=MUTED, interligne=1.4)

    fond = d.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95),
                              Inches(bas - 0.1), Inches(5.4), Inches(2.35))
    fond.fill.solid()
    fond.fill.fore_color.rgb = WHITE
    fond.line.color.rgb = RULE
    fond.line.width = Pt(0.75)
    fond.shadow.inherit = False
    c = _texte(d, 7.35, bas + 0.35, 4.6, 2.4)
    _ligne(c, "« Les candidatures incomplètes ou demeurées incomplètes à la "
              "suite d'une demande de compléments sont éliminées. »",
           taille=16, couleur=INK, police=SERIF, interligne=1.45, apres=12,
           premiere=True)
    _ligne(c, "RÈGLEMENT DE LA CONSULTATION — ANTAI, ARTICLE IV.9",
           taille=9.5, couleur=STAMP, espacement=1.2)

    c = _texte(d, MARGE, 5.75, COLONNE, 1.1)
    _ligne(c, "The buyer has that checklist. The bidder does not.",
           taille=20, gras=True, couleur=INK, premiere=True, apres=8)
    _ligne(c, "A small firm without a bid office is assembling forty documents "
              "across several open tenders, from a folder nobody has re-read "
              "since the last one.",
           taille=14, couleur=MUTED)
    _notes(d, """
Thirty-five seconds. The quotation is the strongest thing on the slide because
it is not your claim — it is the buyer's own wording, in the file itself.
""")

    # 3 ─ what kind of agent ----------------------------------------------
    d = _page(prs)
    bas = _titre(d, "A Professional Agent")
    c = _texte(d, MARGE, bas, 11.4, 1.8)
    _ligne(c, "It does one job inside somebody's working day: the person "
              "assembling a bid, before they submit it.",
           taille=19, couleur=INK, interligne=1.35, premiere=True, apres=10)
    _ligne(c, "Not an assistant, and it does not chat. It reads a document, "
              "checks claims against it, and produces a compliance matrix — the "
              "artefact a bid manager would build by hand over an afternoon, if "
              "they had one.",
           taille=14, couleur=MUTED, interligne=1.4)

    bloc = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGE), Inches(4.45),
                              Inches(COLONNE), Inches(1.55))
    bloc.fill.solid()
    bloc.fill.fore_color.rgb = INK
    bloc.line.fill.background()
    bloc.shadow.inherit = False
    c = _texte(d, MARGE, 4.9, COLONNE, 0.9)
    _ligne(c, "The model observes.  The code decides.", taille=30, gras=True,
           couleur=WHITE, premiere=True, aligne=PP_ALIGN.CENTER)
    c = _texte(d, MARGE, 6.25, COLONNE, 0.6)
    _ligne(c, "Everything after this slide is a proof of that one sentence.",
           taille=13, couleur=MUTED, premiere=True, aligne=PP_ALIGN.CENTER)
    _notes(d, """
Twenty seconds. Say the rule slowly. It is the thesis of the whole submission
and every later beat exists to prove it.
""")

    # 4 ─ what it does -----------------------------------------------------
    d = _page(prs)
    bas = _titre(d, "Four verdicts, because there are four actions")
    verdicts = [
        ("covered", "a document in your library answers it", GOOD),
        ("missing", "nothing does", STAMP),
        ("expires too soon", "it answers today, and not on the submission date",
         STAMP),
        ("to review", "a human has to look", REVIEW),
    ]
    haut = bas
    for etiquette, glose, couleur in verdicts:
        c = _texte(d, MARGE, haut, 3.1, 0.45)
        _ligne(c, etiquette.upper(), taille=12.5, gras=True, couleur=couleur,
               premiere=True, espacement=1.1)
        c = _texte(d, MARGE + 3.35, haut - 0.03, 7.9, 0.5)
        _ligne(c, glose, taille=15, couleur=INK, premiere=True)
        trait = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGE),
                                   Inches(haut + 0.44), Inches(COLONNE),
                                   Inches(0.01))
        trait.fill.solid()
        trait.fill.fore_color.rgb = RULE
        trait.line.fill.background()
        trait.shadow.inherit = False
        haut += 0.72

    c = _texte(d, MARGE, haut + 0.35, COLONNE, 1.2)
    _ligne(c, "And two counts rather than one.", taille=18, gras=True,
           couleur=INK, premiere=True, apres=8)
    _ligne(c, "An incomplete candidature is eliminated. An irregular offre may "
              "be invited to correct itself. The same missing paper ends the "
              "bid in one case and is recoverable in the other — so one number "
              "would tell the reader to treat them alike, which is wrong in "
              "both directions.",
           taille=14, couleur=MUTED, interligne=1.4)
    _notes(d, """
Twenty seconds. Do not read the table aloud line by line — it reads itself.
Spend the words on the two piles, which is the part nobody expects.
""")

    # 5 ─ the documents ----------------------------------------------------
    d = _page(prs)
    bas = _titre(d, "Real files, and one honest fabrication")
    fichiers = [
        ("rc_ANTAI_2026.pdf", "French · 34 pages",
         "Ministry of the Interior. 27 pages store part of their text as images."),
        ("rc_2026SDCRH05.pdf", "French · 14 pages",
         "DGAC. The negative control: nothing hidden."),
        ("itt_EP_COMM_2026.pdf", "English · 16 pages",
         "European Parliament. A turnover floor of EUR 175 000."),
        ("itt_EFSA_2023.pdf", "English · 5 pages",
         "EFSA — and one requirement of it is in French."),
    ]
    haut = bas
    for nom, meta, note in fichiers:
        c = _texte(d, MARGE, haut, 4.0, 0.4)
        _ligne(c, nom, taille=14.5, gras=True, couleur=INK, premiere=True)
        c = _texte(d, MARGE + 4.15, haut + 0.02, 2.0, 0.4)
        _ligne(c, meta, taille=11.5, couleur=STAMP, premiere=True)
        c = _texte(d, MARGE + 6.2, haut + 0.01, 5.2, 0.5)
        _ligne(c, note, taille=12.5, couleur=MUTED, premiere=True)
        haut += 0.62

    c = _texte(d, MARGE, haut + 0.3, COLONNE, 1.6)
    _ligne(c, "All four are published by public buyers and downloadable by "
              "anyone, without registration.",
           taille=14.5, couleur=INK, premiere=True, apres=14)
    _ligne(c, "The evidence library is fabricated, and the repository says so. "
              "Publishing which of a real company's certificates have lapsed is "
              "not something a demonstration gets to do.",
           taille=14.5, gras=True, couleur=STAMP, interligne=1.4)
    _notes(d, """
Twenty-five seconds. The last sentence buys more credibility than anything else
in the video. Say it, and do not hurry it.
""")

    # 6 ─ demo marker: French ---------------------------------------------
    d = _page(prs, INK)
    _filet(d, MARGE, 2.9, largeur=1.6)
    c = _texte(d, MARGE, 3.3, COLONNE, 2.0)
    _ligne(c, "▶  Live run — DGAC, French", taille=40, gras=True, couleur=WHITE,
           premiere=True, apres=16)
    _ligne(c, "17 s on 2026-09-04. Nothing here is replayed.", taille=17,
           couleur=PALE, police=SERIF)
    c = _texte(d, MARGE, 6.2, COLONNE, 0.5)
    _ligne(c, "CUT THIS SLIDE IN THE EDIT — IT MARKS WHERE THE SCREEN "
              "RECORDING GOES", taille=10.5, couleur=MUTED, premiere=True,
           espacement=1.4)
    _notes(d, """
python -X utf8 -m tender_compliance samples/real_dce/rc_2026SDCRH05.pdf \\
    --pages 5-6 --today 2026-08-23 --html out/dgac.html

Talk over the run. Then open out/dgac.html, stop talking, and let them read the
'-9 d' row. Then one sentence: the requirement stays French because the tender
is French and the code checks the quotation against the page it cites — the
English sits beside it, never instead of it.
""")

    # 7 ─ demo marker: English --------------------------------------------
    d = _page(prs, INK)
    _filet(d, MARGE, 2.9, largeur=1.6)
    c = _texte(d, MARGE, 3.3, COLONNE, 2.0)
    _ligne(c, "▶  Live run — European Parliament, English", taille=40,
           gras=True, couleur=WHITE, premiere=True, apres=16)
    _ligne(c, "Same tool, a tender written in English. Watch for what is not "
              "on the screen.", taille=17, couleur=PALE, police=SERIF)
    c = _texte(d, MARGE, 6.2, COLONNE, 0.5)
    _ligne(c, "CUT THIS SLIDE IN THE EDIT — IT MARKS WHERE THE SCREEN "
              "RECORDING GOES", taille=10.5, couleur=MUTED, premiere=True,
           espacement=1.4)
    _notes(d, """
python -X utf8 -m tender_compliance samples/real_dce/itt_EP_COMM_2026.pdf \\
    --today 2026-08-23

35 s on 2026-09-04, which gave 39 requirements, 8 covered, 31 fatal. READ THE
BANNER OFF THE SCREEN — the count moves on every run. Then point at what is NOT
there: no translation lines anywhere. The tool worked out that this document
does not need any.
""")

    # 8 ─ what it found ----------------------------------------------------
    d = _page(prs)
    bas = _titre(d, "What the real files gave back")
    trouvailles = [
        ("−9 d",
         "An insurance certificate valid today and expired on the day bids are "
         "due. Nobody catches that by reading — it is a subtraction."),
        ("27 / 34",
         "Pages of the ANTAI file storing part of their text as images. A "
         "mandatory declaration, legible on screen, invisible to every "
         "extractor tested. So the tool refuses to call anything absent, and "
         "names the pages to open by hand."),
        ("1 of 4",
         "Requirements in the “English” EFSA pack that are written in French. "
         "Found on the first run of a real document — which is why the language "
         "is decided per requirement, not per file."),
    ]
    haut = bas
    for chiffre, texte in trouvailles:
        c = _texte(d, MARGE, haut, 2.3, 0.75)
        _ligne(c, chiffre, taille=34, gras=True, couleur=STAMP, premiere=True)
        c = _texte(d, MARGE + 2.5, haut + 0.05, 8.85, 1.3)
        _ligne(c, texte, taille=14.5, couleur=INK, interligne=1.4, premiere=True)
        haut += 1.55

    c = _texte(d, MARGE, 6.55, COLONNE, 0.6)
    _ligne(c, "Every count moves between runs. Every quotation, and every "
              "verdict computed from it, does not.",
           taille=13, couleur=MUTED, premiere=True)
    _notes(d, """
Twenty-five seconds. If ANTAI is not in your take, drop its line and keep the
other two. The EFSA one is the most surprising and the cheapest to say.
""")

    # 9 ─ how it works -----------------------------------------------------
    d = _page(prs)
    _filet(d, MARGE, 0.62, couleur=STAMP)
    c = _texte(d, MARGE, 0.9, COLONNE, 0.6)
    _ligne(c, "Nothing crosses without being checked", taille=27, gras=True,
           couleur=INK, premiere=True)
    image, rapport = _rendre_le_schema()
    # Fit to the height that is actually left, then centre. Sizing by width and
    # hoping is what pushed the diagram's footnotes off the bottom edge.
    disponible = HAUTEUR - 1.75 - 0.45
    largeur = min(12.4, disponible * rapport)
    d.shapes.add_picture(str(image), Inches((LARGEUR - largeur) / 2), Inches(1.75),
                         width=Inches(largeur))
    _notes(d, """
Thirty-five seconds, over the diagram.

"Two bands. The model proposes on top, twice — the obligations in the text, then
which document might answer each one. It decides nothing. Everything below is
deterministic. A quote that is not on the page it cites is rejected, with a
reason. A document that is not in the library, verbatim, is not a match. Dates
and thresholds are arithmetic and never reach the model at all."
""")

    # 10 ─ Strands ---------------------------------------------------------
    d = _page(prs)
    bas = _titre(d, "Built with the Strands Agents SDK")
    c = _texte(d, MARGE, bas, 11.4, 1.4)
    _ligne(c, "Two agents, and four tools they are given:", taille=17,
           couleur=INK, premiere=True, apres=14)
    outils = [
        ("@tool  page_text", "read a page of the tender"),
        ("@tool  quote_is_on_page", "check a wording really appears on it"),
        ("@tool  list_documents", "list the evidence library"),
        ("@tool  document_is_in_library", "check a name is in it"),
    ]
    haut = bas + 0.75
    for nom, glose in outils:
        c = _texte(d, MARGE + 0.25, haut, 4.3, 0.4)
        _ligne(c, nom, taille=14, gras=True, couleur=INK, police="Consolas",
               premiere=True)
        c = _texte(d, MARGE + 4.9, haut + 0.02, 6.4, 0.4)
        _ligne(c, glose, taille=13.5, couleur=MUTED, premiere=True)
        haut += 0.52

    bloc = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGE), Inches(5.1),
                              Inches(COLONNE), Inches(1.75))
    bloc.fill.solid()
    bloc.fill.fore_color.rgb = INK
    bloc.line.fill.background()
    bloc.shadow.inherit = False
    c = _texte(d, MARGE + 0.5, 5.42, COLONNE - 1.0, 1.2)
    _ligne(c, "They are the same checks that run afterwards — so the agent can "
              "correct a citation before it commits to one.",
           taille=15.5, couleur=PALE, premiere=True, apres=10,
           interligne=1.35)
    _ligne(c, "But it cannot skip them: a tool it never calls changes nothing, "
              "because the verification runs either way.",
           taille=15.5, gras=True, couleur=WHITE, interligne=1.35)
    _notes(d, """
This slide answers judging criterion 1 directly. If you have ten seconds spare,
add the line that shows the discipline is real rather than claimed:

"The English translation is the one thing on the page a model wrote and nothing
verified — so it is the one thing that decides nothing."
""")

    # 11 ─ who it is for ---------------------------------------------------
    d = _page(prs)
    bas = _titre(d, "Who it is for")
    c = _texte(d, MARGE, bas, 11.4, 2.2)
    _ligne(c, "The person assembling the bid: a small firm without a bid "
              "office, a subcontractor, a first-time bidder.",
           taille=20, couleur=INK, interligne=1.35, premiere=True, apres=14)
    _ligne(c, "The buyer already has this checklist. The bidder does not. "
              "Getting it wrong costs a contract that was winnable on the "
              "merits — which is why the tool would rather say “check this” "
              "than say “covered”.",
           taille=15, couleur=MUTED, interligne=1.45)

    bloc = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGE), Inches(4.75),
                              Inches(0.045), Inches(1.5))
    bloc.fill.solid()
    bloc.fill.fore_color.rgb = STAMP
    bloc.line.fill.background()
    bloc.shadow.inherit = False
    c = _texte(d, MARGE + 0.45, 4.85, 10.9, 1.4)
    _ligne(c, "It does not write your bid, and it does not tell you that you "
              "are compliant.", taille=19, gras=True, couleur=INK,
           premiere=True, apres=8, interligne=1.3)
    _ligne(c, "It finds the gaps and shows you where to look.", taille=19,
           couleur=INK, interligne=1.3)
    _notes(d, """
Twenty seconds. Saying plainly what it does not do reads as confidence, not as
a caveat. Do not soften it.
""")

    # 12 ─ what comes next -------------------------------------------------
    d = _page(prs)
    bas = _titre(d, "Recorded, and deliberately not built")
    suite = [
        "One requirement with several satisfaction paths — « ou, à défaut… » — "
        "matching the alternatives rather than only flagging them.",
        "Requirements the buyer can obtain itself from an official register, "
        "and which are legitimately absent from any folder.",
        "Groups of operators: the document checklist multiplies per member "
        "while the capacity thresholds are assessed on the group as a whole.",
    ]
    haut = bas
    for element in suite:
        point = d.shapes.add_shape(MSO_SHAPE.OVAL, Inches(MARGE),
                                   Inches(haut + 0.16), Inches(0.1), Inches(0.1))
        point.fill.solid()
        point.fill.fore_color.rgb = STAMP
        point.line.fill.background()
        point.shadow.inherit = False
        c = _texte(d, MARGE + 0.42, haut, 10.9, 1.1)
        _ligne(c, element, taille=15.5, couleur=INK, interligne=1.4,
               premiere=True)
        haut += 1.05

    c = _texte(d, MARGE, haut + 0.35, COLONNE, 1.0)
    _ligne(c, "Each one is in samples/real_requirements.json with the sentence "
              "that revealed it. Not one was visible from the specification — "
              "they came out of reading real documents.",
           taille=14, couleur=MUTED, interligne=1.4, premiere=True)
    _notes(d, """
Fifteen seconds. This is the first slide to cut if you are over time — it is in
the Devpost text and in the repository.
""")

    # 13 ─ close -----------------------------------------------------------
    d = _page(prs, INK)
    _filet(d, MARGE, 2.35, largeur=1.6)
    c = _texte(d, MARGE, 2.75, COLONNE, 2.2)
    _ligne(c, "Public bids are rejected on paperwork", taille=34, gras=True,
           couleur=WHITE, premiere=True, apres=4, interligne=1.2)
    _ligne(c, "before anyone reads them.", taille=34, gras=True, couleur=WHITE,
           apres=18, interligne=1.2)
    _ligne(c, "This agent finds the gaps first.", taille=26, couleur=PALE,
           police=SERIF)
    c = _texte(d, MARGE, 5.5, COLONNE, 0.8)
    _ligne(c, "github.com/Yoh5/tender-compliance-agent", taille=22, gras=True,
           couleur=WHITE, police="Consolas", premiere=True)
    _notes(d, """
Ten seconds. Leave the URL on screen for a beat after you stop talking, then
end. Do not add a thank-you slide.
""")

    return prs


def main() -> int:
    prs = construire()
    SORTIE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(SORTIE)
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides → {SORTIE}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
